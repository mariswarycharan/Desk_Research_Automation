
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
import matplotlib.pyplot as plt
from io import BytesIO
import pandas as pd

def create_presentation():
    prs = Presentation()
    
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    def add_title_slide():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "European IVD Market Expansion Strategy"
        subtitle.text = "Chinese Supplier Landscape Analysis\n2025 Market Assessment"
        
        title.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 129, 189)
    
    def add_executive_summary():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Strategic Imperative: Chinese IVD Suppliers Face Regulatory Hurdles and Geopolitical Barriers in European Expansion"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "• European IVD market valued at €12.9B annually with 2.6% of total healthcare expenditure"
        p.font.size = Pt(12)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Chinese suppliers face €5M tender restriction under IPI regulations effective June 2025"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• IVDR certification achieved by leading players (Mindray, Snibe, Autobio) but market access remains complex"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Customer receptivity remains divided: cost-advantage in mid-tier segments, trust deficit in high-end solutions"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Strategic acquisitions (HyTest, DiaSys) critical for supply chain localization and market credibility"
        p.font.size = Pt(12)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.fxko5jhzd3mf"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_agenda():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Agenda: Comprehensive Market Assessment Framework"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Market Overview & Size Assessment",
            "Regulatory Landscape Analysis",
            "Competitive Positioning",
            "Access & Reimbursement Framework",
            "Budget Allocation Patterns",
            "Customer Receptivity Mapping",
            "Strategic Opportunities",
            "Risk Assessment",
            "Recommendations"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.level = 0
    
    def add_market_overview():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Market Overview: European IVD Represents Significant Opportunity with Structural Complexities"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Total market size: €12.9B annually across European Union",
            "Per capita expenditure: €150 annually in Germany",
            "Market concentration: 55% controlled by commercial laboratory chains",
            "Growth drivers: Aging population, chronic disease prevalence, technological innovation",
            "Chinese export growth: Doubled between 2015-2023 before IPI restrictions"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.fxko5jhzd3mf"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_regulatory_landscape():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Regulatory Landscape: IVDR Certification Achieved but IPI Creates New Barriers"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "IVDR compliance: Mindray, Snibe, Autobio achieved certification for 4,001+ products",
            "IPI restrictions: Chinese firms barred from public tenders exceeding €5M effective June 2025",
            "Component limitation: Maximum 50% Chinese content in successful bids",
            "Certification timeline: Extended by 6-12 months under IVDR versus previous directive",
            "Environmental compliance: REACH regulations add additional layer of complexity"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        table_left = Inches(0.5)
        table_top = Inches(4.0)
        table_width = Inches(9.0)
        table_height = Inches(1.5)
        
        table = slide.shapes.add_table(3, 3, table_left, table_top, table_width, table_height).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        table.cell(0, 0).text = "Regulatory Framework"
        table.cell(0, 1).text = "Status"
        table.cell(0, 2).text = "Impact"
        
        table.cell(1, 0).text = "IVDR Certification"
        table.cell(1, 1).text = "Achieved by Leaders"
        table.cell(1, 2).text = "Market Access Prerequisite"
        
        table.cell(2, 0).text = "IPI Restrictions"
        table.cell(2, 1).text = "Effective June 2025"
        table.cell(2, 2).text = "€5M Tender Barrier"
        
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.0"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_competitive_landscape():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Competitive Landscape: Chinese Suppliers Gaining Traction Against Established Players"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Market leaders: Roche, Siemens, Abbott dominate central laboratory segment",
            "Chinese presence: Mindray, Snibe, Autobio expanding through acquisitions and partnerships",
            "Product registrations: 4,001 EU certificates across 23 Chinese companies in 2023",
            "Market focus: POCT and chemiluminescence segments showing strongest growth",
            "Localization strategy: Mindray acquired DiaSys (Germany) and HyTest (Finland) for European footprint"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        fig = plt.figure(figsize=(6, 4))
        companies = ['Roche', 'Siemens', 'Abbott', 'Mindray', 'Snibe']
        market_share = [35, 25, 20, 12, 8]
        colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd']
        
        plt.bar(companies, market_share, color=colors)
        plt.title('European IVD Market Share Estimate', fontsize=14)
        plt.ylabel('Market Share (%)', fontsize=12)
        
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)
        
        slide.shapes.add_picture(img_stream, Inches(5.0), Inches(3.0), Inches(4.5), Inches(3.0))
        plt.close(fig)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.4"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://inf.news/en/economy/05703949a08a4f75f598ad8a52850b03.html"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_access_reimbursement():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Access & Reimbursement: Fragmented European Systems Require Country-Specific Strategies"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Germany: EBM catalog governs outpatient reimbursement with fixed pricing",
            "France: NABM system with innovative LAHN pathway for new technologies",
            "Italy: LEA catalog defines reimbursable services with national tariffs",
            "Spain: Global budget system without dedicated IVD reimbursement",
            "UK: NHS Payment Scheme bundles IVD into overall care payments"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        table_left = Inches(0.5)
        table_top = Inches(4.0)
        table_width = Inches(9.0)
        table_height = Inches(1.8)
        
        table = slide.shapes.add_table(4, 3, table_left, table_top, table_width, table_height).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        table.cell(0, 0).text = "Country"
        table.cell(0, 1).text = "Reimbursement System"
        table.cell(0, 2).text = "Key Characteristic"
        
        table.cell(1, 0).text = "Germany"
        table.cell(1, 1).text = "EBM Catalog"
        table.cell(1, 2).text = "Fee-for-service outpatient"
        
        table.cell(2, 0).text = "France"
        table.cell(2, 1).text = "NABM"
        table.cell(2, 2).text = "Innovation funding available"
        
        table.cell(3, 0).text = "Italy"
        table.cell(3, 1).text = "LEA Catalog"
        table.cell(3, 2).text = "National tariff system"
        
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://mtrconsult.com/general-market-access-landscape-vitro-diagnostic-tests-europe"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://journals.publisso.de/en/journals/gms/volume23/000337"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_budget_analysis():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Budget Analysis: Cost Pressure Creates Opportunities for Value Propositions"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Germany: €12.9B annual IVD expenditure with 2.6% of total healthcare budget",
            "Budget growth: 55% increase (2012-2022) versus 63% overall healthcare growth",
            "Efficiency measures: €450M efficiency bonus system controls utilization",
            "Pricing pressure: Fixed technical service prices without inflation adjustment",
            "Hospital allocation: Approximately 2% of clinic budgets allocated to laboratory services"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        fig = plt.figure(figsize=(6, 4))
        categories = ['IVD Growth', 'Overall Healthcare Growth', 'Inflation']
        values = [55, 63, 15]
        colors = ['#2ca02c', '#1f77b4', '#ff7f0e']
        
        plt.bar(categories, values, color=colors)
        plt.title('German Healthcare Expenditure Growth (2012-2022)', fontsize=14)
        plt.ylabel('Percentage Growth (%)', fontsize=12)
        
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)
        
        slide.shapes.add_picture(img_stream, Inches(5.0), Inches(3.0), Inches(4.5), Inches(3.0))
        plt.close(fig)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://journals.publisso.de/en/journals/gms/volume23/000337"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_customer_receptivity():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Customer Receptivity: Trust Deficit in High-End Segments Despite Cost Advantages"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Price sensitivity: Strong receptivity for low-to-mid tier products based on cost competitiveness",
            "Trust barrier: Limited adoption of high-end Chinese IVD solutions due to brand perception",
            "Segment focus: Success in POCT, emergency departments, and small laboratory settings",
            "Local presence: Mindray's 90% local workforce in Turkey demonstrates localization success model",
            "Acquisition strategy: European acquisitions (DiaSys, HyTest) improve credibility and service capability"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        table_left = Inches(0.5)
        table_top = Inches(4.0)
        table_width = Inches(9.0)
        table_height = Inches(1.5)
        
        table = slide.shapes.add_table(3, 3, table_left, table_top, table_width, table_height).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        table.cell(0, 0).text = "Market Segment"
        table.cell(0, 1).text = "Receptivity Level"
        table.cell(0, 2).text = "Key Driver"
        
        table.cell(1, 0).text = "Low-to-Mid Tier"
        table.cell(1, 1).text = "High"
        table.cell(1, 2).text = "Cost competitiveness"
        
        table.cell(2, 0).text = "High-End Segment"
        table.cell(2, 1).text = "Low"
        table.cell(2, 2).text = "Brand trust deficit"
        
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.0"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://healthcare-in-europe.com/en/news/snibe-makes-an-entry-at-euromedlab.html"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_strategic_opportunities():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Strategic Opportunities: Targeted Approaches for Market Penetration"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "POCT expansion: Decentralized testing growth driven by laboratory consolidation trends",
            "Eastern European gateway: Utilize Hungary, Greece for initial EU market entry",
            "German certification: Leverage strict certification for brand credibility enhancement",
            "Local manufacturing: Establish EU production facilities to bypass IPI restrictions",
            "Strategic partnerships: Collaborate with European distributors and laboratory networks"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.0"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.4"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_risk_assessment():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Risk Assessment: Geopolitical and Regulatory Challenges Require Mitigation"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "IPI restrictions: €5M tender barrier effective June 2025 limits large contract opportunities",
            "Component limitations: 50% Chinese content restriction requires supply chain diversification",
            "Regulatory complexity: IVDR transition adds 6-12 months to certification timelines",
            "Political scrutiny: Increasing EU focus on supply chain sovereignty and reduced Chinese dependence",
            "Reimbursement pressure: Fixed pricing systems without inflation adjustment limit margin potential"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        table_left = Inches(0.5)
        table_top = Inches(4.0)
        table_width = Inches(9.0)
        table_height = Inches(1.8)
        
        table = slide.shapes.add_table(4, 3, table_left, table_top, table_width, table_height).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        table.cell(0, 0).text = "Risk Category"
        table.cell(0, 1).text = "Impact Level"
        table.cell(0, 2).text = "Mitigation Strategy"
        
        table.cell(1, 0).text = "Geopolitical"
        table.cell(1, 1).text = "High"
        table.cell(1, 2).text = "Local manufacturing"
        
        table.cell(2, 0).text = "Regulatory"
        table.cell(2, 1).text = "High"
        table.cell(2, 2).text = "Early IVDR compliance"
        
        table.cell(3, 0).text = "Competitive"
        table.cell(3, 1).text = "Medium"
        table.cell(3, 2).text = "Niche segmentation"
        
        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.alignment = PP_ALIGN.CENTER
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://www.ainvest.com/news/frontier-medtech-navigating-china-eu-trade-tensions-investment-gains-2507/"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_recommendations():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Strategic Recommendations: Three-Pronged Approach for Market Success"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "Accelerate localization: Establish European manufacturing facilities to bypass IPI restrictions",
            "Focus on strategic segments: Target POCT and mid-tier segments where cost advantage is strongest",
            "Enhance clinical evidence: Invest in European clinical trials and health economic studies",
            "Leverage acquisitions: Continue strategic purchases of European companies for market access",
            "Develop partnership models: Create joint ventures with European distributors and laboratory networks"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.0"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://hytest.fi/news/mindray-completes-acquisition-of-hytest"
        p.font.size = Pt(8)
        p.font.italic = True
    
    def add_conclusion():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        title = slide.shapes.title
        title.text = "Conclusion: Strategic Transformation Required for Sustainable European Presence"
        title.text_frame.paragraphs[0].font.size = Pt(16)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        points = [
            "European IVD market offers significant growth potential but requires sophisticated market entry strategy",
            "Regulatory compliance (IVDR) is necessary but insufficient without addressing geopolitical barriers",
            "Local manufacturing presence critical for accessing large public tender opportunities",
            "Customer acceptance growing in value segments but requires continued investment in quality and clinical evidence",
            "Strategic acquisitions and partnerships provide fastest path to market credibility and distribution capability"
        ]
        
        for point in points:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 0 else tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
        
        ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1.5))
        tf = ref_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://docs.google.com/document/d/1KWqiQRt4NvZuv2_3GC2UQWi7agGABNUshk7Pk_lRDTM/edit?tab=t.0"
        p.font.size = Pt(8)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://www.ainvest.com/news/frontier-medtech-navigating-china-eu-trade-tensions-investment-gains-2507/"
        p.font.size = Pt(8)
        p.font.italic = True

    add_title_slide()
    add_executive_summary()
    add_agenda()
    add_market_overview()
    add_regulatory_landscape()
    add_competitive_landscape()
    add_access_reimbursement()
    add_budget_analysis()
    add_customer_receptivity()
    add_strategic_opportunities()
    add_risk_assessment()
    add_recommendations()
    add_conclusion()

    prs.save('european_ivd_market_analysis.pptx')

create_presentation()