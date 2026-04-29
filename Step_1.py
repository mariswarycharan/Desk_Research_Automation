import asyncio
import io
import json
import os
import re
import sys

import pandas as pd
import pdfplumber
import requests
from crawl4ai import AsyncWebCrawler
from docx import Document
from pptx import Presentation
from ollama import Client

MODEL = "deepseek-v3.1:671b-cloud"
client = Client(
    host="https://ollama.com",
    headers={'Authorization': 'Bearer dd456319486541c5bfd8dfd001136b32.krWKOjU-1cHaHoKQr0iQDe0r'}
)

OUTPUT_FILE = "Research_Evidence_Report.xlsx"

# ─── STEP 1: CONTENT EXTRACTION ───────────────────────────────────────────────

def extract_from_pdf(url_or_path: str) -> str | None:
    """Extract text from a PDF via URL or local path."""
    try:
        if url_or_path.startswith("http"):
            resp = requests.get(url_or_path, headers={"User-Agent": "Mozilla/5.0"}, timeout=15, verify=False)
            if resp.status_code != 200:
                return None
            source = io.BytesIO(resp.content)
        else:
            source = url_or_path

        with pdfplumber.open(source) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except Exception as e:
        print(f"   [PDF Error] {e}")
        return None


def extract_from_docx(path: str) -> str | None:
    """Extract text from a .docx file (local path only)."""
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        print(f"   [DOCX Error] {e}")
        return None


def extract_from_pptx(path: str) -> str | None:
    """Extract text from a .pptx file (local path only)."""
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"   [PPTX Error] {e}")
        return None


def extract_from_txt(path: str) -> str | None:
    """Read plain text file."""
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"   [TXT Error] {e}")
        return None


async def extract_from_website(url: str) -> str | None:
    """Scrape web content using Crawl4AI."""
    try:
        async with AsyncWebCrawler(verbose=False) as crawler:
            result = await crawler.arun(url=url)
            return result.markdown
    except Exception as e:
        print(f"   [Web Error] {e}")
        return None


def detect_and_extract_local(path: str) -> str | None:
    """Route local file to the correct extractor based on extension."""
    ext = os.path.splitext(path)[1].lower()
    extractors = {
        ".pdf": extract_from_pdf,
        ".docx": extract_from_docx,
        ".pptx": extract_from_pptx,
        ".txt": extract_from_txt,
    }
    fn = extractors.get(ext)
    if fn:
        return fn(path)
    print(f"   [Unsupported format] {ext}")
    return None


async def extract_content(source: str) -> str | None:
    """
    Main dispatcher: detects source type and returns extracted text.
    Handles: website URL, PDF URL, or local file (pdf/docx/pptx/txt).
    """
    if source.startswith("http"):
        if source.lower().endswith(".pdf"):
            return extract_from_pdf(source)
        return await extract_from_website(source)
    else:
        return detect_and_extract_local(source)


# ─── STEP 2: DYNAMIC LLM EXTRACTION ──────────────────────────────────────────

# topic declaration
topic_1 = "Access_Reimbursement"
topic_1_description = """
Extract information related to market access, reimbursement policies, insurance coverage, public or private payer support, pricing approvals, and patient affordability. Include evidence on funding pathways, inclusion in healthcare schemes, tender access, and barriers to diagnosis adoption due to cost or reimbursement limitations.
"""

topic_2 = "Budget_Spend"
topic_2_description = """
Extract information related to healthcare spending, diagnostic budgets, procurement investments, hospital or laboratory expenditure, government allocations, and funding trends or costs , money, Include data on purchasing capacity, capital investment, cost pressures, and budget priorities influencing IVD adoption.
"""

topic_3 = "Regulatory"
topic_3_description = """
Extract information related to regulatory approvals, compliance requirements, certification pathways, import/export rules, quality standards, and market authorization processes. Include updates on IVDR, CE marking, local regulations, policy changes, and barriers impacting market entry.
"""

topic_4 = "Customer_Receptivity"
topic_4_description = """
Extract information related to customer acceptance, clinician adoption, laboratory demand, distributor interest, brand perception, purchasing behavior, and stakeholder readiness. Include evidence on preferences for new technologies, trust in suppliers, unmet needs, and willingness to switch or adopt Chinese IVD products.
"""


def build_prompt(text: str, research_topic: str, research_domain: str, focus_entities: str) -> str:
    """
    Builds a dynamic system+user prompt suited for any research context.
    All research parameters are injected at runtime.
    """
    safe_text = text

    return f"""
You are a Senior Research Intelligence Analyst. Your task is to extract structured evidence from a source document for the following research context:

- Research Domain  : {research_domain}
- Research Topic   : {research_topic}
- Focus Entities   : {focus_entities}

-----

INPUT SOURCE CONTENT (Extracted text from a document or webpage):
{safe_text}

-----


Your extraction must answer these five analytical pillars where applicable:
1. {topic_1}: {topic_1_description} \n
2. {topic_2}: {topic_2_description} \n
3. {topic_3}: {topic_3_description} \n
4. {topic_4}: {topic_4_description} \n
5. Any other relevant insights \n

---

IMPORTANT INSTRUCTIONS:

1. Do not be overly strict while extracting relevant information for each topic based on its description. If any content in the source is even slightly relevant to a topic, include it.\n

2. Understand the context and meaning of each sentence carefully before extracting the information. Capture what the content is actually trying to convey.\n

3. If multiple sentences or pieces of information are relevant to a particular topic, include all of them. Present them clearly as bullet points such as Point 1, Point 2, etc.\n

4. Ensure no relevant information is missed while extracting the data. \n


OUTPUT RULES:
- Return ONLY a valid JSON object. No markdown, no preamble.
- For every pillar: include Status (Yes/No), the EXACT first sentence of the supporting statement from the source, and a strategic Explanation.
- Include a "Summary" of 5-8 sentences synthesizing the full source relevance to the research topic.

JSON STRUCTURE:
{{
"Source Title": "Give suitable title of the source document or webpage",
"Name of the Source": "Give suitable name of the publication or website",
"Source Type": "Give suitable Report, News, Policy, Market Data, etc.",
"Year of Publication": "give me Year if available",
"Period of the Data Used": "give me e.g. 2018-2023, 2024 or 'Not specified'",
"Description": "Brief short description of the source content",
"Access Type": "Free, Subscription, Paywall, Open Access, etc.",
"Source Link": "URL if applicable",
"Geography Focus": "France/Italy/Germany/Spain/UK/Global/etc.",
"Companies Mentioned (Focus)": "Comma-separated list of focus entities mentioned in the source",
"Other Entities Mentioned": "Comma-separated list of any other notable organizations/products mentioned",

"Entities_Found": "Comma-separated focus entities found (or 'None')",
"Other_Entities": "Other notable organizations/products mentioned",

"Access_Reimbursement": {{
"Status": "Yes or No",
"Exact_Opening_Sentence": "First sentence of the verbatim evidence",
"Summarized_Evidence": "Paraphrased synthesis of the full evidence block",
"Explanation": "Strategic implication for the research topic"
}},

"Budget_Spend": {{
"Status": "Yes or No",
"Exact_Opening_Sentence": "...",
"Summarized_Evidence": "...",
"Explanation": "..."
}},

"Regulatory": {{
"Status": "Yes or No",
"Exact_Opening_Sentence": "...",
"Summarized_Evidence": "...",
"Explanation": "..."
}},

"Customer_Receptivity": {{
"Status": "Yes or No",
"Exact_Opening_Sentence": "...",
"Summarized_Evidence": "...",
"Explanation": "..."
}},

"Others": {{
"Exact_Opening_Sentence": "...",
"Summarized_Evidence": "...",
"Explanation": "..."
}},

"Summary": "4-5 sentence synthesis of source relevance to the research."
}}
"""

async def analyze_with_llm(text: str, research_topic: str, research_domain: str, focus_entities: str) -> dict | None:
    """Send extracted text to LLaMA via Ollama and return parsed JSON."""
    prompt = build_prompt(text, research_topic, research_domain, focus_entities)

    for attempt in range(3):
        try:
            messages = [
                {"role": "system", "content": "You are a research extraction engine. Output valid JSON only."},
                {"role": "user", "content": prompt}
            ]

            # Ollama client.chat is synchronous, run in executor to avoid blocking
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(
                None,
                lambda: client.chat(MODEL, messages=messages)
            )

            raw = response['message']['content'].strip()

            # save raw response for reference as txt file
            with open(f"raw_response.txt", "w", encoding="utf-8") as f:
                f.write(raw)

            # Strip markdown fences if present
            raw = re.sub(r"^```json|^```|```$", "", raw, flags=re.MULTILINE).strip()
            return json.loads(raw)

        except json.JSONDecodeError as e:
            print(f"   [JSON Parse Error] {e}")
            return None
        except Exception as e:
            if "429" in str(e) or "rate" in str(e).lower():
                print(f"   [Rate Limit] Retrying in 60s... (attempt {attempt + 1})")
                await asyncio.sleep(60)
            else:
                print(f"   [LLM Error] {e}")
                return None
    return None

# ─── STEP 3: FLATTEN JSON → EXCEL ROW ────────────────────────────────────────

def flatten_to_row(source: str, data: dict) -> dict:
    """Map LLM JSON output to a flat Excel row dictionary."""

    def pillar(key):
        block = data.get(key, {})
        return {
            "status": block.get("Status", "No"),
            "opening": block.get("Exact_Opening_Sentence", ""),
            "summary": block.get("Summarized_Evidence", ""),
            "explanation": block.get("Explanation", ""),
        }

    ar = pillar("Access_Reimbursement")
    bs = pillar("Budget_Spend")
    rg = pillar("Regulatory")
    cr = pillar("Customer_Receptivity")
    ot = pillar("Others")

    return {
        "Source": source,
        "Source Title": data.get("Source Title", ""),
        "Name of the Source": data.get("Name of the Source", ""),
        "Source Type": data.get("Source Type", ""),
        "Year of Publication": data.get("Year of Publication", ""), 
        "Period of the Data Used": data.get("Period of the Data Used", ""),
        "Description": data.get("Description", ""),
        "Access Type": data.get("Access Type", ""),
        "Source Link": data.get("Source Link", ""),
        "Geography Focus": data.get("Geography Focus", ""),
        "Companies Mentioned (Focus)": data.get("Companies Mentioned (Focus)", ""),
        "Other Entities Mentioned": data.get("Other Entities Mentioned", ""),

        "Entities Found (Focus)": data.get("Entities_Found", "None"),
        "Other Entities Mentioned": data.get("Other_Entities", "None"),

        "Access & Reimbursement (Yes/No)": ar["status"],
        "Access – Exact Opening Sentence": ar["opening"],
        "Access – Summarized Evidence": ar["summary"],
        "Access – Explanation": ar["explanation"],

        "Budget & Spend (Yes/No)": bs["status"],
        "Budget – Exact Opening Sentence": bs["opening"],
        "Budget – Summarized Evidence": bs["summary"],
        "Budget – Explanation": bs["explanation"],

        "Regulatory (Yes/No)": rg["status"],
        "Regulatory – Exact Opening Sentence": rg["opening"],
        "Regulatory – Summarized Evidence": rg["summary"],
        "Regulatory – Explanation": rg["explanation"],

        "Customer Receptivity (Yes/No)": cr["status"],
        "Customer – Exact Opening Sentence": cr["opening"],
        "Customer – Summarized Evidence": cr["summary"],
        "Customer – Explanation": cr["explanation"],

        "Others – Exact Opening Sentence": ot["opening"],
        "Others – Summarized Evidence": ot["summary"],
        "Others – Explanation": ot["explanation"],

        "Overall Summary": data.get("Summary", ""),
    }


# ─── MAIN PIPELINE ────────────────────────────────────────────────────────────

async def run_pipeline(
    sources: list[str],
    research_topic: str,
    research_domain: str,
    focus_entities: str,
    output_file: str = OUTPUT_FILE,
):
    """
    Full pipeline: extract → analyze → export.

    Args:
        sources        : List of URLs or local file paths.
        research_topic : e.g. "European IVD market entry strategy"
        research_domain: e.g. "In-Vitro Diagnostics (IVD)"
        focus_entities : e.g. "Mindray, Snibe"
        output_file    : Path to save the Excel report.
    """
    print(f"\n=== Research Extraction Pipeline ===")
    print(f"Topic   : {research_topic}")
    print(f"Domain  : {research_domain}")
    print(f"Entities: {focus_entities}")
    print(f"Sources : {len(sources)}\n")

    rows = []

    for i, source in enumerate(sources, 1):
        print(f"[{i}/{len(sources)}] {source}")

        # Step 1: Extract
        text = await extract_content(source)
        if not text or len(text.strip()) < 200:
            print("   > Skipped (insufficient content)")
            continue
        print(f"   > Extracted {len(text):,} characters")

        # Step 2: Analyze
        data = await analyze_with_llm(text, research_topic, research_domain, focus_entities)
        if not data:
            print("   > Skipped (LLM returned no data)")
            continue

        # Step 3: Flatten
        row = flatten_to_row(source, data)
        rows.append(row)
        print("   > Row added to report")

    if not rows:
        print("\nNo data extracted. Exiting.")
        return

    df = pd.DataFrame(rows)
    df.to_excel(output_file, index=False)
    print(f"\n=== Done. Report saved to: {output_file} ===")


# ─── ENTRY POINT ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    # ── CONFIGURE YOUR RUN HERE ──────────────────────────────────────────────
    SOURCES = [
        # Mix of URLs and local files — all supported
        "https://healthcare-in-europe.com/en/news/snibe-makes-an-entry-at-euromedlab.html#:~:text=%E2%80%98Our%20exports%20to%20Europe%20focus,It%E2%80%99s%20perfectly%C2%A0suited%2C%20for%20example%2C%20for",

    ]

    RESEARCH_TOPIC = "European expansion potential of Chinese IVD suppliers"
    RESEARCH_DOMAIN = "In-Vitro Diagnostics (IVD)"
    FOCUS_ENTITIES = "Mindray, Snibe"
    # ─────────────────────────────────────────────────────────────────────────

    # # just call website extraction to test the extraction without calling the LLM
    # web_content = asyncio.run(extract_from_website(SOURCES[0]))

    # # just call only build prompt to test the prompt output without calling the LLM
    # result = build_prompt(
    #     text=web_content,
    #     research_topic=RESEARCH_TOPIC,
    #     research_domain=RESEARCH_DOMAIN,
    #     focus_entities=FOCUS_ENTITIES
    # )
    # print("Generated Prompt:")
    # print(result)

    asyncio.run(run_pipeline(
        sources=SOURCES,
        research_topic=RESEARCH_TOPIC,
        research_domain=RESEARCH_DOMAIN,
        focus_entities=FOCUS_ENTITIES,
        output_file=OUTPUT_FILE,
    ))