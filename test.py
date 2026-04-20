
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib.pyplot as plt
from io import BytesIO
import pandas as pd

def create_presentation():
    prs = Presentation()
    
    slide_layouts = prs.slide_layouts
    title_slide_layout = slide_layouts[0]
    title_only_layout = slide_layouts[5]
    blank_layout = slide_layouts[6]
    content_layout = slide_layouts[1]
    
    blue_color = RGBColor(0, 32, 96)
    
    def add_title_slide():
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "European IVD Market Analysis: Chinese Supplier Expansion"
        subtitle.text = "Strategic Research Briefing\n2025 Market Intelligence Report"
    
    def add_executive_summary():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Executive Summary: Chinese IVD Suppliers Face High-Value European Market with Significant Barriers"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "• European IVD market represents €1.65 trillion healthcare expenditure with €3,685 per capita spending"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Chinese suppliers face new EU procurement restrictions banning participation in tenders >€5M with >50% Chinese inputs"
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = "• IVDR transition creates significant regulatory barrier, lengthening certification by 6-12 months and forcing 30% of SMEs out"
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = "• Mindray leads Chinese expansion with 3,914 global registrations (3,158 overseas) and acquisition of German DiaSys"
        p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = "• Cost-effectiveness drives adoption in low-mid segments but trust deficit hampers high-end market penetration"
        p.font.size = Pt(14)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[4] https://en.caclp.com/industry-news/2365.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_agenda():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Agenda"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "1. Market Overview & Key Metrics"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "2. Regulatory Landscape Analysis"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "3. Competitive Intelligence"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "4. Market Access & Procurement Barriers"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "5. Strategic Recommendations"
        p.font.size = Pt(16)
        p.font.bold = True
    
    def add_market_overview():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "European IVD Market: High-Value Opportunity with Complex Structure"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Market Size & Structure"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• €1.65 trillion total healthcare expenditure with €3,685 per capita spending"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Government and mandatory insurance funding represents over 60% of budget structure"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Tiered reimbursement system: DRG-based prospective payment at base level, supplementary mechanisms for innovative technologies"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Key Market Characteristics"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Not monolithic - collection of government-funded systems with public tender rules"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Cost-control paramount but mechanisms exist for premium reimbursement of proven clinical value"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Price sensitivity toward low-to-mid-end devices in developed European markets"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_regulatory_landscape():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Regulatory Landscape: IVDR Creates Significant Barriers for Market Entry"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "IVDR Transition Impact"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Certification process lengthened by 6-12 months since MDR implementation"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• More than 30% of small and medium-sized enterprises forced out of market"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Transition periods extended to 2027-2028 requiring strict compliance throughout product lifecycle"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Chinese Company Progress"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Snibe: 211 chemiluminescence + 67 biochemical reagents with IVDR CE certification"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Mindray: First biochemical project IVDR CE certificate (β2-microglobulin) in January 2021"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Wondfo: First IVDR quality management system certificate in POCT field (August 2021)"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Over 40 Chinese companies granted CE for monkeypox virus detection kits (2022)"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[7] https://en.caclp.com/industry-news/3530.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[8] https://en.caclp.com/industry-news/1160.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[10] https://en.caclp.com/industry-news/1480.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def create_regulatory_chart():
        fig, ax = plt.subplots(figsize=(8, 4))
        companies = ['Snibe', 'Mindray', 'Wondfo', 'Other Chinese\nCompanies']
        certifications = [278, 1, 1, 40]
        colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728']
        
        bars = ax.bar(companies, certifications, color=colors)
        ax.set_ylabel('Number of IVDR Certifications', fontweight='bold')
        ax.set_title('Chinese IVD Companies: IVDR Certification Progress', fontweight='bold')
        
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 5,
                    f'{int(height)}', ha='center', va='bottom', fontweight='bold')
        
        plt.tight_layout()
        
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=300)
        img_stream.seek(0)
        plt.close()
        
        return img_stream
    
    def add_regulatory_chart_slide():
        slide = prs.slides.add_slide(blank_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Chinese IVD Suppliers Making Significant Progress on IVDR Compliance"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = blue_color
        
        img_stream = create_regulatory_chart()
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[7] https://en.caclp.com/industry-news/3530.html | [8] https://en.caclp.com/industry-news/1160.html | [10] https://en.caclp.com/industry-news/1480.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_procurement_barriers():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Procurement Restrictions: EU Limits Chinese Participation in High-Value Tenders"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "EU International Procurement Instrument (IPI) Measures"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Exclusion from public procurement tenders for medical devices >€5 million contract value"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Maximum 50% Chinese inputs allowed for successful bids"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Exceptions only where no alternative suppliers exist"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Market Impact Assessment"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Targets high-value public tenders central to IVD lab consolidation deals"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Chinese suppliers limited to smaller, decentralized budget allocations"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Non-Chinese companies may face issues if supply chain relies heavily on Chinese OEM components"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Retaliatory measure against China's 87% exclusion rate of EU medical devices"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def create_procurement_table():
        slide = prs.slides.add_slide(blank_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "EU Procurement Restrictions: Impact Analysis by Contract Value"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = blue_color
        
        rows = 4
        cols = 3
        left = Inches(1.5)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(2)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(2.5)
        
        table.cell(0, 0).text = "Contract Value Threshold"
        table.cell(0, 1).text = "Chinese Participation"
        table.cell(0, 2).text = "Market Impact"
        
        table.cell(1, 0).text = "Above €5 million"
        table.cell(1, 1).text = "Excluded"
        table.cell(1, 2).text = "Blocked from large-scale lab consolidation tenders"
        
        table.cell(2, 0).text = "Below €5 million"
        table.cell(2, 1).text = "Permitted"
        table.cell(2, 2).text = "Access to smaller, decentralized purchases"
        
        table.cell(3, 0).text = "Any value with >50% Chinese inputs"
        table.cell(3, 1).text = "Restricted"
        table.cell(3, 2).text = "Supply chain diversification required"
        
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row, col)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.bold = (row == 0)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_competitive_landscape():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Competitive Landscape: Chinese IVD Companies Expanding Global Presence"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Market Position & Global Reach"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Mindray: 3,914 global registrations (3,158 overseas) - leader in international expansion"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Snibe: 1,493 global registration certificates with strong overseas instrument sales"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Wondfo: Global service network covering 150+ countries with localized production bases"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Strategic Acquisitions & Partnerships"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Mindray acquired Germany's DiaSys (2023) and Finland's HyTest Invest Oy (2021)"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Autobio partnership with Sekisui Medical for biochemical reagents in China"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Wondfo acquired Shenzhen Tisenc Medical to expand chemiluminescence business"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[3] https://en.caclp.com/industry-news/1161.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[4] https://en.caclp.com/industry-news/2365.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[5] https://en.caclp.com/industry-news/2785.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def create_competitive_table():
        slide = prs.slides.add_slide(blank_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Chinese IVD Leaders: Global Registration Footprint and Strategic Positioning"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = blue_color
        
        rows = 4
        cols = 4
        left = Inches(0.8)
        top = Inches(1.3)
        width = Inches(8.5)
        height = Inches(2.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(2)
        table.columns[3].width = Inches(2)
        
        table.cell(0, 0).text = "Company"
        table.cell(0, 1).text = "Global Registrations"
        table.cell(0, 2).text = "Overseas Registrations"
        table.cell(0, 3).text = "Key Strategic Moves"
        
        table.cell(1, 0).text = "Mindray"
        table.cell(1, 1).text = "3,914"
        table.cell(1, 2).text = "3,158"
        table.cell(1, 3).text = "Acquired DiaSys (DE), HyTest (FI); 30% int'l growth"
        
        table.cell(2, 0).text = "Snibe"
        table.cell(2, 1).text = "1,493"
        table.cell(2, 2).text = "N/A"
        table.cell(2, 3).text = "278 IVDR products; strong chemiluminescence focus"
        
        table.cell(3, 0).text = "Wondfo"
        table.cell(3, 1).text = "N/A"
        table.cell(3, 2).text = "N/A"
        table.cell(3, 3).text = "150+ countries; acquired Tisenc; localized production"
        
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row, col)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.bold = (row == 0)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link | [4] https://en.caclp.com/industry-news/2365.html | [7] https://en.caclp.com/industry-news/3530.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_financial_performance():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Financial Performance: Post-COVID Adjustment with Strong Export Growth"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Market Contraction & Recovery"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• 2023: 53.26% revenue decline to CNY 108.44B from CNY 231.99B in 2022 across 59 IVD companies"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• 84.5% net profit decline to CNY 9.14B from CNY 59.12B in 2022"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Only 17 of 59 firms showed year-on-year revenue growth in 2023"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Leading Performers & Export Growth"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Mindray: 6 consecutive years of >20% profit growth; 30% international IVD business growth"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Overseas revenue growth: Mindray +21.28%, Autobio +27.6%"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• R&D investment: Mindray allocated 10.99% of operating income (RMB 1.030B)"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[5] https://en.caclp.com/industry-news/2785.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[6] https://en.caclp.com/industry-news/2808.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[9] https://en.caclp.com/industry-news/3349.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def create_revenue_chart():
        years = ['2022', '2023', '2024']
        total_revenue = [231.99, 108.44, 139.28]
        mindray_revenue = [0, 0, 13.76]
        
        fig, ax = plt.subplots(figsize=(8, 5))
        x = range(len(years))
        
        width = 0.35
        bars1 = ax.bar([i - width/2 for i in x], total_revenue, width, label='Total Industry (59 Companies)', color='#1f77b4')
        bars2 = ax.bar([i + width/2 for i in x], mindray_revenue, width, label='Mindray IVD Division', color='#ff7f0e')
        
        ax.set_xlabel('Year', fontweight='bold')
        ax.set_ylabel('Revenue (CNY Billion)', fontweight='bold')
        ax.set_title('Chinese IVD Market Revenue Trends (2022-2024)', fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(years)
        ax.legend()
        
        for bars in [bars1, bars2]:
            for bar in bars:
                height = bar.get_height()
                if height > 0:
                    ax.text(bar.get_x() + bar.get_width()/2., height + 5,
                            f'{height}', ha='center', va='bottom', fontsize=9, fontweight='bold')
        
        plt.tight_layout()
        
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=300)
        img_stream.seek(0)
        plt.close()
        
        return img_stream
    
    def add_revenue_chart_slide():
        slide = prs.slides.add_slide(blank_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Chinese IVD Market: Post-COVID Revenue Correction with Leading Players Recovering"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = blue_color
        
        img_stream = create_revenue_chart()
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[6] https://en.caclp.com/industry-news/2808.html | [9] https://en.caclp.com/industry-news/3349.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_customer_receptivity():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Customer Receptivity: Cost-Effectiveness vs. Trust Deficit in High-End Market"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Adoption Drivers"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Price competitiveness drives adoption in low-to-mid-tier segments"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Chinese exports to EU more than doubled between 2015-2023 indicating historical receptivity"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Snibe sold 3,637 chemiluminescent instruments abroad vs 1,141 domestically (2022)"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Challenges & Barriers"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Low brand recognition in developed European markets"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Trust deficit for high-end products from Chinese manufacturers"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Market dominated by international giants (Roche, J&J) with strong brands and advanced technologies"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Need for local KOL endorsement and price incentives to displace existing workflows"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[4] https://en.caclp.com/industry-news/2365.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_strategic_approaches():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Strategic Approaches: Market Entry Strategies for European Expansion"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Market Entry Pathways"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Eastern European gateway countries (Hungary, Greece) for initial EU market access"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• German/Dutch certification for enhanced brand credibility despite stricter regulations"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Targeting small laboratories to middle-size diagnostic centers with compact POCT solutions"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Localization Strategies"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Acquisition of European entities (Mindray: DiaSys Germany, HyTest Finland)"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Establishing local service networks and production capabilities"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Partnerships with European companies for market access and credibility"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Product Segmentation"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Focus on high-margin chemiluminescence sector (20-25% localization rate in China)"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Development of 'world's first' technical specifications for high-end market appeal"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Emergency response kits (monkeypox) as gateway products for market entry"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[3] https://en.caclp.com/industry-news/1161.html"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[10] https://en.caclp.com/industry-news/1480.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_risks_challenges():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Risks & Challenges: Navigating Complex European Market Dynamics"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Regulatory & Compliance Risks"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• IVDR transition lengthening certification by 6-12 months with increased costs"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Environmental regulations (REACH) adding compliance complexity for reagents and instruments"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Need for continuous compliance throughout product lifecycle with extended transition to 2027-2028"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Market Access Barriers"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• EU procurement restrictions blocking participation in tenders >€5M with >50% Chinese inputs"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• EU push for supply chain localization increasing uncertainty for pure importers"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Fragmented reimbursement systems across EU member states requiring country-specific strategies"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Competitive Challenges"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Trust deficit in high-end segments dominated by Roche, Siemens, Abbott"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Need for significant local KOL endorsement to displace established workflows"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Technical gap compared to global leaders in registration certificate footprint"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[4] https://en.caclp.com/industry-news/2365.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_recommendations():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Strategic Recommendations: Navigating the European IVD Market Landscape"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Market Entry Strategy"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Focus on sub-€5M tenders and private sector deals to bypass procurement restrictions"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Utilize Eastern European countries (Hungary, Greece) as gateway markets for EU access"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Pursue German certification for enhanced brand credibility despite regulatory complexity"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Operational Approach"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Acquire local European entities for established distribution and service capabilities"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Establish European manufacturing footprint to mitigate supply chain localization policies"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Develop local KOL relationships to build trust and overcome brand recognition deficit"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Product & Market Focus"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Target high-margin chemiluminescence and POCT segments with compact solutions"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Leverage emergency response capabilities for rapid market entry during health crises"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Balance cost-effectiveness with innovation to compete beyond price-sensitive segments"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[7] https://en.caclp.com/industry-news/3530.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    def add_conclusion():
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = "Conclusion: Strategic Imperatives for European IVD Market Success"
        title.text_frame.paragraphs[0].font.color.rgb = blue_color
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Key Takeaways"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• European IVD market offers high-value opportunity but requires navigating complex barriers"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• IVDR compliance is table stakes but procurement restrictions create additional hurdles"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Chinese suppliers must transition from export model to localized presence for long-term success"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Future Outlook"
        p.font.size = Pt(14)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• Market shifting from simple export economics to deep localization requirements"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Companies with European acquisitions and local capabilities will gain competitive advantage"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Regulatory compliance and procurement restrictions may evolve based on trade reciprocity"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "• Innovation and quality improvements essential to overcome trust deficit in high-end segments"
        p.font.size = Pt(12)
        
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "[1] https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[2] https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
        p.font.size = Pt(10)
        p.font.italic = True
        
        p = tf.add_paragraph()
        p.text = "[4] https://en.caclp.com/industry-news/2365.html"
        p.font.size = Pt(10)
        p.font.italic = True
    
    add_title_slide()
    add_executive_summary()
    add_agenda()
    add_market_overview()
    add_regulatory_landscape()
    add_regulatory_chart_slide()
    add_procurement_barriers()
    create_procurement_table()
    add_competitive_landscape()
    create_competitive_table()
    add_financial_performance()
    add_revenue_chart_slide()
    add_customer_receptivity()
    add_strategic_approaches()
    add_risks_challenges()
    add_recommendations()
    add_conclusion()
    
    prs.save('European_IVD_Market_Analysis_Chinese_Expansion.pptx')

create_presentation()