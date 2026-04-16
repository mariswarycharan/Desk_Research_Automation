from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
import io

prs = Presentation()

def add_title_slide():
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "European IVD Market Analysis: Chinese Supplier Expansion"
    subtitle.text = "Strategic Research Presentation\nClient Confidential"

def add_slide_with_title(title_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    return slide

def add_bullet_slide(title_text, content):
    slide = add_slide_with_title(title_text)
    text_box = slide.shapes.placeholders[1]
    tf = text_box.text_frame
    tf.text = ""
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
    return slide

def add_references(slide, refs):
    if not refs:
        return
    ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
    tf = ref_box.text_frame
    tf.text = "References:"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True
    for i, ref in enumerate(refs, 1):
        p = tf.add_paragraph()
        p.text = f"[{i}] {ref}"
        p.level = 0
        p.font.size = Pt(9)
        p.font.italic = True

def create_company_table(slide, data):
    rows, cols = len(data) + 1, len(data[0])
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(3)
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
    for col_idx, header in enumerate(data[0]):
        table.cell(0, col_idx).text = header
        table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
    for row_idx, row_data in enumerate(data[1:], 1):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx, col_idx).text = str(cell_data)
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_regulatory_timeline_chart():
    fig, ax = plt.subplots(figsize=(10, 4))
    events = [
        ('IVDR Implementation Start', '2022-05-26'),
        ('Snibe First IVDR Cert', '2020-12-01'),
        ('Mindray First IVDR Cert', '2021-01-01'),
        ('Wondfo IVDR QMS Cert', '2021-08-01'),
        ('Snibe 18 Reagents IVDR', '2025-08-20')
    ]
    dates = [pd.to_datetime(date) for _, date in events]
    labels = [label for label, _ in events]
    ax.plot(dates, [1]*len(dates), marker='o', linestyle='-', markersize=8)
    ax.set_yticks([])
    for i, (label, date) in enumerate(events):
        ax.annotate(label, (pd.to_datetime(date), 1), xytext=(5, 15),
                   textcoords='offset points', rotation=45, ha='left')
    ax.set_title('Key IVDR Regulatory Milestones for Chinese Companies')
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def create_market_growth_chart():
    fig, ax = plt.subplots(figsize=(8, 5))
    companies = ['Mindray', 'Snibe', 'Autobio']
    growth_rates = [30, 16.62, 27.6]
    colors = ['#1f77b4', '#ff7f0e', '#2ca02c']
    bars = ax.bar(companies, growth_rates, color=colors)
    ax.set_ylabel('Overseas Revenue Growth (%)')
    ax.set_title('Chinese IVD Company Overseas Growth (2023-2024)')
    for bar in bars:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height + 0.5,
                f'{height}%', ha='center', va='bottom')
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300)
    plt.close()
    img_stream.seek(0)
    return img_stream

def add_image_to_slide(slide, img_stream, left, top, width, height):
    slide.shapes.add_picture(img_stream, left, top, width, height)

add_title_slide()

exec_summary_content = [
    "European IVD market represents high-value opportunity with €1.65T healthcare expenditure",
    "Chinese suppliers face significant regulatory barriers with IVDR transition extending certification by 6-12 months",
    "Geopolitical factors creating new challenges with EU restricting Chinese participation in contracts >€5M",   
    "Leading Chinese companies (Mindray, Snibe, Autobio) achieving IVDR certification and showing strong overseas growth",
    "Strategic acquisitions (Mindray/DiaSys, Mindray/HyTest) key to overcoming trust deficit and supply chain risks",
    "Cost-effectiveness drives adoption in low-mid segments, but high-end market remains dominated by Western incumbents"
]
exec_slide = add_bullet_slide("Executive Summary", exec_summary_content)
add_references(exec_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
])

agenda_content = [
    "Market Overview & Size",
    "Key Trends & Developments",
    "Competitive Landscape",
    "Regulatory Environment",
    "Access & Reimbursement",
    "Budget & Spend Analysis",
    "Customer Receptivity",
    "Strategic Opportunities",
    "Risks & Challenges",
    "Recommendations"
]
agenda_slide = add_bullet_slide("Agenda", agenda_content)

market_content = [
    "European healthcare expenditure: €1.65 trillion with €3,685 per capita spending",
    "Government and mandatory insurance funding represents over 60% of budget structure",
    "Market transitioning from COVID-era volumes back to traditional non-COVID business",
    "Chinese IVD sector showing contraction with total revenue dropping 53.26% in 2023",
    "Long-term demand drivers: population aging, chronic disease rates, health awareness"
]
market_slide = add_bullet_slide("Market Overview", market_content)
add_references(market_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://en.caclp.com/industry-news/2808.html"
])

trends_content = [
    "IVDR transition creating market consolidation with 30% of SMEs forced out",
    "Geopolitical tensions leading to EU restrictions on Chinese participation in large contracts",
    "Strategic acquisitions accelerating (Mindray/DiaSys, Mindray/HyTest, Wondfo/Tisenc)",
    "Localization strategies becoming critical for market access and customer trust",
    "R&D investment increasing with Mindray allocating 10.99% of revenue to medical R&D",
    "Shift toward AI and automation integration in product development"
]
trends_slide = add_bullet_slide("Key Trends & Developments", trends_content)
add_references(trends_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://en.caclp.com/industry-news/2785.html",
    "https://en.caclp.com/industry-news/3349.html"
])

competitive_data = [
    ["Company", "Global Registrations", "Overseas Registrations", "2024 Growth", "IVDR Progress"],
    ["Mindray", "3,914", "3,158", "30%+", "Multiple certifications"],
    ["Snibe", "1,493", "N/A", "16.62%", "211 chemiluminescence reagents"],
    ["Maccura", "2,467", "N/A", "N/A", "N/A"],
    ["Autobio", "N/A", "N/A", "27.6%", "IVDR certified"],
    ["Wondfo", "N/A", "140+ countries", "N/A", "IVDR QMS certified"]
]
competitive_slide = add_slide_with_title("Competitive Landscape")
create_company_table(competitive_slide, competitive_data)
add_references(competitive_slide, [
    "https://en.caclp.com/industry-news/2365.html",
    "https://en.caclp.com/industry-news/2785.html"
])

regulatory_content = [
    "IVDR (EU 2017/746) implementation began May 26, 2022 with extended transition to 2027-2028",
    "CE certification mandatory for market access, acting as regulatory 'visa'",
    "Chinese companies among first to achieve IVDR certifications across multiple categories",
    "REACH environmental regulations adding compliance complexity for reagents and instruments",
    "High-risk IVD products require Notified Body certification; low-risk can be self-declared",
    "Average certification process lengthened by 6-12 months under IVDR"
]
regulatory_slide = add_bullet_slide("Regulatory Landscape", regulatory_content)
timeline_img = create_regulatory_timeline_chart()
add_image_to_slide(regulatory_slide, timeline_img, Inches(1), Inches(3.5), Inches(8), Inches(3))
add_references(regulatory_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://en.caclp.com/industry-news/1160.html",
    "https://en.caclp.com/industry-news/3530.html"
])

reimbursement_content = [
    "EU member states use pricing systems driven by public health insurance",
    "National/regional agencies negotiate with manufacturers balancing affordability, profits, compliance",       
    "Tiered reimbursement structure: DRG-based prospective payment system plus supplementary mechanisms",
    "Public tender rules and long-term cost control paramount in government-funded systems",
    "European market fragmentation requires country-specific reimbursement strategies"
]
reimbursement_slide = add_bullet_slide("Access & Reimbursement", reimbursement_content)
add_references(reimbursement_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
])

budget_content = [
    "€5M threshold for EU restrictions on Chinese participation in public procurement",
    "Chinese suppliers limited to smaller, decentralized budget allocations",
    "Successful bids must contain no more than 50% Chinese inputs",
    "Exceptions possible where no alternative suppliers exist",
    "Government and mandatory insurance represents majority of funding (>60%)"
]
budget_slide = add_bullet_slide("Budget & Spend Analysis", budget_content)
add_references(budget_slide, [
    "https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569",
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link"
])

receptivity_content = [
    "Chinese medical device exports to EU more than doubled between 2015-2023",
    "Price sensitivity in low-to-mid-end devices, brand trust deficit in high-end products",
    "Snibe sold 3,637 units abroad vs 1,141 domestically in first three quarters of 2022",
    "Market dominated by international giants (Roche, J&J, Siemens, Abbott) with strong brands",
    "Chinese companies targeting small to middle size diagnostic centers with compact solutions",
    "Performance-to-cost ratio creating substitution effect in international markets"
]
receptivity_slide = add_bullet_slide("Customer Receptivity", receptivity_content)
growth_img = create_market_growth_chart()
add_image_to_slide(receptivity_slide, growth_img, Inches(4), Inches(3.5), Inches(5), Inches(3.5))
add_references(receptivity_slide, [
    "https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569",
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://en.caclp.com/industry-news/2365.html"
])

opportunities_content = [
    "Strategic acquisitions of European companies to secure supply chains and credibility",
    "Targeting Eastern European countries (Hungary, Greece) as gateways to EU market",
    "Leveraging German/Dutch certification for enhanced brand credibility",
    "Focus on high-margin chemiluminescence sector with 20-25% localization rate in China",
    "Developing compact POCT solutions for decentralized testing segments",
    "Utilizing cost-effectiveness advantage for price-sensitive market segments"
]
opportunities_slide = add_bullet_slide("Strategic Opportunities", opportunities_content)
add_references(opportunities_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://en.caclp.com/industry-news/1161.html"
])

risks_content = [
    "EU supply chain localization policies increasing uncertainty for Chinese exporters",
    "Complex international regulatory and certification barriers (FDA PMA, EU CE-MDR)",
    "Technology and intellectual property risks in developed markets",
    "Difficulty in localization and building brand trust in high-end segments",
    "Geopolitical tensions affecting market access through IPI restrictions",
    "Market fragmentation requiring country-specific strategies"
]
risks_slide = add_bullet_slide("Risks & Challenges", risks_content)
add_references(risks_slide, [
    "https://drive.google.com/file/d/1zt3vNjOlugypsvPkywduyWNRPGTHrsoj/view?usp=drive_link",
    "https://ec.europa.eu/commission/presscorner/detail/en/ip_25_1569"
])

recommendations_content = [
    "Establish local European manufacturing and service capabilities through acquisitions",
    "Develop country-specific market entry strategies accounting for reimbursement variations",
    "Invest in high-level IVDR compliance and environmental regulation (REACH) expertise",
    "Build strategic partnerships with European distributors and service providers",
    "Focus product development on high-margin segments with competitive advantages",
    "Diversify geographic focus to include Eastern European gateway markets",
    "Develop strong clinical evidence and KOL relationships to overcome trust deficit"
]
recommendations_slide = add_bullet_slide("Recommendations", recommendations_content)

conclusion_content = [
    "European IVD market presents significant opportunity but requires sophisticated approach",
    "Regulatory compliance (IVDR) is foundational but insufficient without localization strategy",
    "Geopolitical factors increasingly influencing market access conditions",
    "Leading Chinese companies demonstrating capability to compete through acquisitions and innovation",
    "Success requires balancing cost-effectiveness with quality perception and local presence",
    "Market shifting from export economics to requirement for deep European localization"
]
conclusion_slide = add_bullet_slide("Conclusion", conclusion_content)

sources_slide = add_slide_with_title("Sources")
sources_content = [
    "Frost & Sullivan: 2025 Blue Book on Global Expansion of Chinese Medical Devices",
    "European Commission: Restrictions on Chinese participation in medical devices procurement",
    "CACLP: Various industry news articles on Chinese IVD company performance",
    "Multiple company announcements and regulatory filings"
]
text_box = sources_slide.shapes.placeholders[1]
tf = text_box.text_frame
tf.text = ""
for source in sources_content:
    p = tf.add_paragraph()
    p.text = source
    p.level = 0
    p.font.size = Pt(14)

prs.save('ddddddddddddEuropean_IVD_Market_Analysis.pptx')